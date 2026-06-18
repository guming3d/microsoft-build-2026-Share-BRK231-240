"""
Generate a PPTX summarizing the WebIQ feature investigation:
 - code-test findings (Q1 freshness, Q2 pagination, Q4 cold-start, Q7 browse)
 - open questions that CANNOT be answered by code and need vendor confirmation
   (Q4 SLA/regions, Q7 crawl legality, Q8 content rights / attribution / compliance)

Run:  python make_findings_ppt.py
Out:  WebIQ_Findings_and_Open_Questions.pptx  (in this folder)
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---- palette ----------------------------------------------------------------
BLUE = RGBColor(0x00, 0x78, 0xD4)     # Microsoft blue
DARK = RGBColor(0x20, 0x20, 0x20)
GREY = RGBColor(0x60, 0x60, 0x60)
GREEN = RGBColor(0x10, 0x7C, 0x10)
RED = RGBColor(0xC4, 0x3E, 0x1C)
AMBER = RGBColor(0xB7, 0x6E, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF3, 0xF6, 0xFB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def rect(s, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def set_run(r, text, size, color=DARK, bold=False, italic=False, font="Segoe UI"):
    r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = font


def header(s, title, kicker=None):
    rect(s, 0, 0, 13.333, 1.15, BLUE)
    _, tf = box(s, 0.55, 0.18, 12.2, 0.85)
    p = tf.paragraphs[0]
    set_run(p.add_run(), title, 28, WHITE, bold=True)
    if kicker:
        p2 = tf.add_paragraph()
        set_run(p2.add_run(), kicker, 13, RGBColor(0xD9, 0xEC, 0xFB))


def bullets(s, items, l, t, w, h, size=15, gap=6):
    _, tf = box(s, l, t, w, h)
    first = True
    for it in items:
        # Two accepted shapes:
        #   (level:int, text:str, [color], [bold])
        #   (text:str, [color], [bold])   -> level 0
        #   "text"                        -> level 0
        if isinstance(it, str):
            lvl, txt, color, bold = 0, it, DARK, False
        elif isinstance(it[0], int):
            lvl = it[0]
            txt = it[1]
            color = it[2] if len(it) > 2 else DARK
            bold = it[3] if len(it) > 3 else False
        else:
            lvl = 0
            txt = it[0]
            color = it[1] if len(it) > 1 else DARK
            bold = it[2] if len(it) > 2 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.level = lvl
        bullet = "\u25aa " if lvl == 0 else ("\u2013 " if lvl == 1 else "\u00b7 ")
        set_run(p.add_run(), bullet + txt, size - lvl, color, bold=bold)


def table(s, rows, l, t, w, col_w, head_color=BLUE, fsize=12):
    nrows, ncols = len(rows), len(rows[0])
    gw = s.shapes.add_table(nrows, ncols, Inches(l), Inches(t), Inches(w), Inches(0.4 * nrows)).table
    for ci, cw in enumerate(col_w):
        gw.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gw.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            txt, color, bold = val, DARK, False
            if isinstance(val, tuple):
                txt, color, bold = val[0], val[1], (val[2] if len(val) > 2 else False)
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = head_color
                set_run(p.add_run(), str(txt), fsize, WHITE, bold=True)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                set_run(p.add_run(), str(txt), fsize, color, bold=bold)
    return gw


def footer(s, n):
    _, tf = box(s, 0.5, 7.02, 12.3, 0.4)
    p = tf.paragraphs[0]
    set_run(p.add_run(),
            "WebIQ (Foundry) limited-access evaluation  \u00b7  verified via live API tests, 2026-06-18",
            9, GREY)
    p.alignment = PP_ALIGN.LEFT
    _, tf2 = box(s, 12.4, 7.02, 0.6, 0.4)
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    set_run(p2.add_run(), str(n), 9, GREY)


# ============================================================ SLIDE 1 - TITLE
s = slide()
rect(s, 0, 0, 13.333, 7.5, BLUE)
rect(s, 0, 5.0, 13.333, 2.5, RGBColor(0x00, 0x5A, 0x9E))
_, tf = box(s, 0.9, 1.7, 11.5, 2.6)
set_run(tf.paragraphs[0].add_run(), "Microsoft WebIQ (Foundry)", 40, WHITE, bold=True)
p = tf.add_paragraph(); set_run(p.add_run(), "Feature & Performance Findings + Open Questions", 26, RGBColor(0xCF, 0xE5, 0xFA))
p = tf.add_paragraph(); p.space_before = Pt(14)
set_run(p.add_run(), "Verified through live API test cases  \u00b7  plus questions requiring vendor confirmation", 15, RGBColor(0xCF, 0xE5, 0xFA))
_, tf = box(s, 0.9, 5.5, 11.5, 1.5)
set_run(tf.paragraphs[0].add_run(), "Scope: Q1 result freshness  \u00b7  Q2 pagination  \u00b7  Q4 performance  \u00b7  Q7 browse  \u00b7  Q8 content rights", 14, WHITE)
p = tf.add_paragraph(); set_run(p.add_run(), "Base URL  https://api.microsoft.ai/v3", 13, RGBColor(0xCF, 0xE5, 0xFA))

# ============================================================ SLIDE 2 - APPROACH
s = slide()
header(s, "Approach: what we tested vs. what we must ask", "Five user-critical questions, split by whether code can answer them")
table(s, [
    ["#", "Question", "Method", "Status"],
    [("Q1", DARK, True), "Result determinism, cache TTL, freshness filter", "Code test (live)", ("Answered", GREEN, True)],
    [("Q2", DARK, True), "Pagination / coverage beyond maxResults", "Code test (live)", ("Answered", GREEN, True)],
    [("Q4", DARK, True), "Latency: cold-start vs warm", "Code test (live)", ("Answered", GREEN, True)],
    [("Q4", DARK, True), "Is 164ms a contractual SLA? Regional endpoints?", "Vendor confirm", ("Open", AMBER, True)],
    [("Q7", DARK, True), "Browse: robots.txt, paywall, JS rendering", "Code test (live)", ("Answered", GREEN, True)],
    [("Q7", DARK, True), "Legal right to crawl/store crawled content", "Vendor confirm", ("Open", AMBER, True)],
    [("Q8", DARK, True), "Caching rights, attribution, GDPR/residency", "Vendor confirm", ("Open", AMBER, True)],
], 0.55, 1.5, 12.2, [0.7, 5.7, 2.6, 1.7], fsize=13)
bullets(s, [
    ("All code tests run against the live tenant with a real API key; each finding below is reproducible via the named script.", GREY, False),
], 0.55, 6.2, 12.2, 0.6, size=12)
footer(s, 2)

# ============================================================ SLIDE 3 - RESULTS SUMMARY
s = slide()
header(s, "Code-test findings at a glance", "Four questions answered with live evidence")
table(s, [
    ["Question", "Finding", "Evidence"],
    [("Q1 Determinism", DARK, True), ("Same query = identical results; no random variation", GREEN, False), "5x calls, 100% Jaccard"],
    [("Q1 Cache", DARK, True), "Result set stable for full 5-min window (TTL >= 300s)", "see TTL probe"],
    [("Q1 Freshness", DARK, True), ("news freshness filter is a no-op (ignored)", RED, True), "raw _post probe"],
    [("Q2 Pagination", DARK, True), ("NONE - maxResults is a hard ceiling", RED, True), "8 paging params ignored"],
    [("Q2 Ranking", DARK, True), "web re-orders by size; news top-10 = prefix of top-20", "subset test"],
    [("Q4 Cold-start", DARK, True), ("7-15x penalty: cold 1.2-1.7s vs warm 0.1-0.2s", RED, True), "5 rounds/endpoint"],
    [("Q7 robots.txt", DARK, True), ("Enforced - LinkedIn/Facebook -> HTTP 403 dropped", GREEN, False), "browse() calls"],
    [("Q7 JS render", DARK, True), "Cached pages ignore render flag; live crawl = 202 + ~4.6s", "uncached probe"],
], 0.55, 1.45, 12.2, [2.5, 6.7, 3.0], fsize=12)
footer(s, 3)

# ============================================================ SLIDE 4 - Q1
s = slide()
header(s, "Q1 - Result freshness & caching", "How stable are results, and can I get fresher ones?")
bullets(s, [
    ("Determinism: 5 repeated calls returned the EXACT same result set (100% overlap) for both a specific and a time-sensitive query.", DARK, True),
    (1, "maxResults controls HOW MANY items, never WHICH items - you always get the same top-N."),
    ("Caching: the first call is cold (~1.2-1.7s); repeats drop to ~0.1-0.2s, indicating a cache layer.", DARK, True),
    (1, "Result set stayed identical across a 5-minute window -> cache TTL is at least 300s."),
    ("Freshness control: the news endpoint does NOT expose a working freshness filter.", DARK, True),
    (1, "day/week/month/year were accepted (HTTP 200) but returned the identical 20 results (no-op)."),
    ("Cache-bust: no noCache / bypassCache / fresh body param changed the response.", DARK, True),
    ("Implication: for changing content you must vary the query itself; there is no client-side 'force fresh'.", BLUE, True),
], 0.55, 1.5, 12.2, 5.2, size=15, gap=7)
footer(s, 4)

# ============================================================ SLIDE 5 - Q2
s = slide()
header(s, "Q2 - Pagination & coverage", "Can I get results beyond the maxResults cap?")
bullets(s, [
    ("There is NO pagination. maxResults is a hard ceiling on what you can ever retrieve for a query.", RED, True),
    (1, "Probed 8 conventional params - offset, skip, page, from, cursor, start, pageNumber, count."),
    (1, "Every one was silently ignored: identical page-1 results, no 'page 2'."),
    ("Per-query maximums: web 50, news 20, images 30, videos 30, classic 50 (web answers).", DARK, True),
    ("Ranking is not size-stable for web search:", DARK, True),
    (1, "top-10 is a SUBSET of top-50 but NOT a strict prefix - order shifts when you change maxResults."),
    (1, "news is better-behaved: top-10 was an exact prefix of top-20."),
    ("Implication: to widen coverage you must reformulate the query (sub-topics, filters); you cannot deep-page.", BLUE, True),
    ("Design tip: always request the max you need in ONE call; do not assume top-10 == first 10 of top-50.", BLUE, True),
], 0.55, 1.5, 12.2, 5.2, size=15, gap=7)
footer(s, 5)

# ============================================================ SLIDE 6 - Q4
s = slide()
header(s, "Q4 - Performance: cold-start vs warm", "What latency should I design for?")
table(s, [
    ["Endpoint", "Cold (1st call)", "Warm (reused conn)", "Penalty"],
    [("web", DARK, True), "1683 ms", "182 ms", ("+1500 ms (9.2x)", RED, True)],
    [("news", DARK, True), "1424 ms", "200 ms", ("+1224 ms (7.1x)", RED, True)],
    [("browse (cached)", DARK, True), "1505 ms", "100 ms", ("+1405 ms (15.1x)", RED, True)],
], 0.55, 1.5, 9.0, [2.6, 2.2, 2.2, 2.0], fsize=13)
bullets(s, [
    ("Connection reuse is essential: a fresh TCP/TLS handshake adds 1.2-1.5s every time.", DARK, True),
    (1, "Use a pooled, keep-alive HTTP session (the client reuses requests.Session) for all calls."),
    ("Warm latency (100-200ms) is realistic for steady traffic; the documented 164ms p95 is only reachable warm.", DARK, True),
    (1, "Only browse-on-cache met 164ms; live-search endpoints sit ~180-250ms warm."),
    ("Uncached browse live-crawl is much slower: returns 202 then takes ~4.6s to complete.", DARK, True),
], 0.55, 3.5, 12.2, 3.2, size=14, gap=6)
footer(s, 6)

# ============================================================ SLIDE 7 - Q7
s = slide()
header(s, "Q7 - Browse behavior", "What does /browse actually return for hard pages?")
bullets(s, [
    ("robots.txt IS respected / sites can be blocked:", DARK, True),
    (1, "LinkedIn and Facebook both returned HTTP 403 'result is dropped' - no content served."),
    ("Paywalled sites return thin content:", DARK, True),
    (1, "WSJ section page -> only ~273 chars (preview/metadata); NYT -> ~5,471 chars."),
    (1, "You should NOT assume full article text behind a hard paywall."),
    ("JavaScript rendering (renderDynamicPages):", DARK, True),
    (1, "Only engages on a LIVE crawl. Cached pages ignore the flag (identical output, ~95ms)."),
    (1, "Uncached pages return HTTP 202 'crawl in progress'; a full live crawl + render took ~4.6s."),
    (1, "renderDynamicPages requires liveCrawl='fallback' (validated)."),
    ("Implication: design browse for async behavior - handle 202 + retry, and budget seconds for live crawls.", BLUE, True),
], 0.55, 1.45, 12.2, 5.4, size=14, gap=5)
footer(s, 7)

# ============================================================ SLIDE 8 - OPEN Q4
s = slide()
header(s, "OPEN \u2013 Q4: SLA & infrastructure", "Cannot be answered by code \u2013 needs vendor confirmation")
rect(s, 0.55, 1.45, 12.2, 0.5, AMBER)
_, tf = box(s, 0.7, 1.5, 12.0, 0.45)
set_run(tf.paragraphs[0].add_run(), "These determine whether WebIQ can carry a production workload with predictable behavior.", 13, WHITE, bold=True)
bullets(s, [
    ("Is the documented 164ms p95 a contractual SLA, or a best-effort target?", DARK, True),
    (1, "We measured 182-250ms warm; need the official latency commitment and the percentile/window it's measured over."),
    ("What is the availability SLA (uptime %), and what are the service credits on breach?", DARK, True),
    ("Are there regional / geo-distributed endpoints to cut latency, or a single global endpoint?", DARK, True),
    (1, "Affects users far from the host region; today only one base URL is documented."),
    ("What are the real rate limits & quotas (RPS, per-minute, daily, monthly)?", DARK, True),
    (1, "We burst ~22 req/s with no 429; the true ceiling and the 429 retryAfter policy are unknown."),
    ("Versioning & deprecation: how much notice before /v3 is retired or changed?", DARK, True),
], 0.55, 2.2, 12.2, 4.6, size=15, gap=8)
footer(s, 8)

# ============================================================ SLIDE 9 - OPEN Q7
s = slide()
header(s, "OPEN \u2013 Q7: Crawl rights & limits", "Legal/operational questions code cannot settle")
rect(s, 0.55, 1.45, 12.2, 0.5, AMBER)
_, tf = box(s, 0.7, 1.5, 12.0, 0.45)
set_run(tf.paragraphs[0].add_run(), "We can observe behavior, but not the contractual terms or exact thresholds behind it.", 13, WHITE, bold=True)
bullets(s, [
    ("What exactly triggers HTTP 430 'too many crawls' - requests per second / minute / hour?", DARK, True),
    (1, "Needed to size live-crawl workloads safely; we did not hit it in light testing."),
    ("Is content from /browse licensed for storage, indexing, and downstream display to my users?", DARK, True),
    (1, "Observed 403 on robots-restricted sites, but the terms of use for permitted content are a contract question."),
    ("For paywalled / metered sites, what is the official position - is partial content expected and permitted to reuse?", DARK, True),
    ("What is the freshness of the browse cache, and can a true real-time fetch be guaranteed?", DARK, True),
    (1, "Cached responses (~90ms) may be stale; only liveCrawl forces a fetch, at multi-second cost."),
    ("SLA / timeout policy for the 202 'crawl in progress' path - max wait before failure?", DARK, True),
], 0.55, 2.2, 12.2, 4.6, size=15, gap=7)
footer(s, 9)

# ============================================================ SLIDE 10 - OPEN Q8
s = slide()
header(s, "OPEN \u2013 Q8: Content rights & compliance", "The highest-risk open area \u2013 entirely non-testable")
rect(s, 0.55, 1.45, 12.2, 0.5, RED)
_, tf = box(s, 0.7, 1.5, 12.0, 0.45)
set_run(tf.paragraphs[0].add_run(), "Blocks production/legal sign-off until answered by Microsoft.", 13, WHITE, bold=True)
bullets(s, [
    ("Caching & storage: may I cache returned results, snippets, URLs and crawled text? For how long?", DARK, True),
    ("Attribution: what attribution / branding is required when displaying results to end users?", DARK, True),
    ("Redistribution: can results be shown to external/third-party users, or internal-only?", DARK, True),
    ("Data residency: where are queries processed and stored? Any in-region (EU) guarantee?", DARK, True),
    ("Privacy/GDPR: are query strings logged, retained, or used for training? Can I opt out / get a DPA?", DARK, True),
    ("PII handling: what happens if a query or crawled page contains personal data?", DARK, True),
    ("Compliance posture: ISO 27001, SOC 2, FedRAMP, EU Data Boundary coverage?", DARK, True),
], 0.55, 2.2, 12.2, 4.6, size=15, gap=10)
footer(s, 10)

# ============================================================ SLIDE 11 - NEXT STEPS
s = slide()
header(s, "Recommendations & next steps", "Turning findings into action")
bullets(s, [
    ("Engineering - act on confirmed findings now:", BLUE, True),
    (1, "Reuse a keep-alive HTTP session everywhere (avoids the 1.2-1.5s cold-start tax)."),
    (1, "Request the full maxResults in one call; never deep-page - it isn't supported."),
    (1, "Handle browse 202 + retry; budget multiple seconds for uncached live crawls."),
    (1, "Don't rely on news 'freshness'; vary the query to refresh content."),
    ("Vendor - get written answers before production sign-off:", BLUE, True),
    (1, "Q4: latency & uptime SLA, regional endpoints, real rate limits/quota, version deprecation policy."),
    (1, "Q7: 430 crawl thresholds, license to store/show crawled content, cache freshness guarantees."),
    (1, "Q8: caching rights, attribution, data residency, GDPR/DPA, compliance certifications. (HIGHEST PRIORITY)"),
    ("Re-test on each release to catch doc/behavior drift (e.g., docs say HTTP->410, live API returns 400).", DARK, True),
], 0.55, 1.5, 12.2, 5.2, size=14, gap=6)
footer(s, 11)

out = "WebIQ_Findings_and_Open_Questions.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
